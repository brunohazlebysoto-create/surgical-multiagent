import base64
import hashlib
import json
import logging
import os
from typing import Optional

import fitz
import httpx
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.agents.base import call_gemini

logger = logging.getLogger("document_parser")


def extract_text_from_pdf(filepath: str) -> str:
    """Extrae texto de un archivo PDF usando pypdf."""
    try:
        reader = PdfReader(filepath)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extrayendo texto de PDF {filepath}: {e}")
        raise e


def extract_text_from_docx(filepath: str) -> str:
    """Extrae texto de un archivo de Word (.docx) usando python-docx."""
    try:
        doc = Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        # También extraer texto de tablas
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_cells = [cell.text for cell in row.cells if cell.text]
                if row_cells:
                    table_text.append(" | ".join(row_cells))
        return "\n".join(paragraphs + table_text).strip()
    except Exception as e:
        logger.error(f"Error extrayendo texto de DOCX {filepath}: {e}")
        raise e


def extract_text_from_pptx(filepath: str) -> str:
    """Extrae texto de una presentación de PowerPoint (.pptx) usando python-pptx."""
    try:
        prs = Presentation(filepath)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n--- Diapositiva {i + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extrayendo texto de PPTX {filepath}: {e}")
        raise e


async def generate_document_abstract(
    text: str, filename: str, pdf_base64: Optional[str] = None
) -> str:
    """
    Toma el texto completo extraído (o el PDF codificado en Base64 para modo multimodal)
    y usa Gemini para generar un resumen estructurado (Abstract) médico/quirúrgico de alta calidad.
    """
    if pdf_base64:
        prompt = f"""
        Eres un transcriptor y redactor médico especializado en cirugía infantil y medicina basada en evidencia.
        Analiza el documento PDF clínico adjunto (Archivo: {filename}).
        El usuario requiere que no te limites al resumen; debes leer detalladamente todo el contenido del documento, analizando minuciosamente cada tabla, gráfico, figura, diagrama e imagen clínica presente en el PDF.

        Por favor, genera un resumen/síntesis estructurado extremadamente riguroso y detallado (Abstract ampliado) de máximo 600 palabras que incluya:
        1. **Objetivo y Diseño**: Objetivo principal, tipo de estudio y el contexto clínico detallado.
        2. **Población Pediátrica**: Detalles específicos de los pacientes (edad promedio/grupos etarios, peso, criterios de inclusión/exclusión y severidad de la patología).
        3. **Detalles Quirúrgicos/Médicos y Figuras/Imágenes**: Describe con precisión técnica las imágenes diagnósticas (radiografías, ecografías, TAC, etc.), diagramas quirúrgicos, pasos de la intervención, y los gráficos de resultados (por ejemplo, curvas de aprendizaje, diagramas de barras, etc.).
        4. **Tablas de Datos y Resultados**: Transcribe los datos y estadísticas clave de las tablas del documento (éxito %, complicaciones %, p-values, odds ratios, etc.).
        5. **Conclusiones Clínicas**: Hallazgos clave aplicables directamente en el quirófano infantil.

        Este resumen enriquecido servirá como la base de conocimiento completa de este paper para un panel de 15 agentes de IA. No omitas datos numéricos ni hallazgos visuales.
        """
        try:
            logger.info(f"Llamando a Gemini en modo MULTIMODAL para el PDF: {filename}")
            summary = await call_gemini(
                prompt,
                temperature=0.15,
                inline_data={"mimeType": "application/pdf", "data": pdf_base64},
            )
            return summary
        except Exception as e:
            logger.error(
                f"Error generando resumen multimodal para {filename}: {e}. Intentando fallback con texto extraído."
            )

    if not text:
        return "El archivo no contiene texto legible."

    # Limitar texto de entrada si es absurdamente largo por seguridad de tokens (ej: 30k caracteres)
    input_text = text[:30000]

    prompt = f"""
    Eres un transcriptor y redactor médico especializado en cirugía infantil y medicina basada en evidencia.
    Analiza el siguiente documento clínico subido por el usuario (Archivo: {filename}).
    El usuario requiere que no te limites al resumen; debes leer detalladamente todo el contenido del texto, extrayendo información de las tablas, gráficos, figuras e imágenes descritas en el texto.

    Texto completo del documento:
    ---
    {input_text}
    ---

    Genera un resumen/síntesis estructurado extremadamente riguroso y detallado (Abstract ampliado) de máximo 600 palabras que incluya:
    1. **Objetivo y Diseño**: Objetivo principal, tipo de estudio y el contexto clínico detallado.
    2. **Población Pediátrica**: Detalles específicos de los pacientes (edad promedio/grupos etarios, peso, criterios de inclusión/exclusión y severidad de la patología).
    3. **Detalles Quirúrgicos/Médicos y Figuras**: Pasos exactos de la intervención quirúrgica (colocación de puertos, instrumental, reparos anatómicos) y descripción detallada de hallazgos en imágenes diagnósticas, gráficos de supervivencia, curvas de aprendizaje o figuras de resultados mencionadas.
    4. **Tablas de Datos y Resultados**: Transcribe los datos clave de las tablas (porcentajes de éxito, tasas de complicaciones estadísticas en %, tiempos quirúrgicos, sangrado estimado, etc.).
    5. **Conclusiones Clínicas**: Hallazgos clave aplicables directamente en el quirófano infantil.

    Este resumen enriquecido será utilizado por un panel de 15 agentes de IA para un meta-análisis GRADE y redacción de consensos de alta precisión. No resumas por encima, conserva los datos numéricos y comparativos exactos.
    """

    try:
        summary = await call_gemini(prompt, temperature=0.15)
        return summary
    except Exception as e:
        logger.error(
            f"Error generando resumen de texto para {filename} con Gemini: {e}"
        )
        # Fallback simple
        return f"Documento subido '{filename}'. Contenido de texto extraído: {text[:500]}..."


async def parse_uploaded_document(
    filepath: str, filename: str, run_id: Optional[str] = None
) -> dict:
    """
    Orquesta la extracción del texto y la generación del abstract/resumen.
    Devuelve un diccionario estructurado como un 'paper'.
    """
    ext = os.path.splitext(filepath)[1].lower()
    text = ""
    pdf_base64 = None
    extracted_imgs = []

    logger.info(f"Procesando archivo subido: {filename} ({ext})")

    if ext == ".pdf":
        # Extraer texto de respaldo por si falla el modo multimodal
        try:
            text = extract_text_from_pdf(filepath)
        except Exception as err:
            logger.warning(f"No se pudo extraer texto de PDF {filename}: {err}")
            text = ""

        # Extraer e interpretar imágenes del PDF si se cuenta con run_id
        if run_id:
            try:
                extracted_imgs = await extract_and_analyze_images(filepath, run_id)
            except Exception as e:
                logger.error(f"Error durante la extracción de imágenes de PDF: {e}")

        # Codificar en Base64 para multimodal si tiene tamaño adecuado (<15MB)
        try:
            file_size = os.path.getsize(filepath)
            if file_size <= 15 * 1024 * 1024:
                with open(filepath, "rb") as f:
                    pdf_bytes = f.read()
                pdf_base64 = base64.b64encode(pdf_bytes).decode("utf-8")
                logger.info(
                    f"PDF {filename} codificado en Base64 ({len(pdf_base64)} caracteres) para análisis multimodal."
                )
            else:
                logger.warning(
                    f"El archivo {filename} ({file_size} bytes) excede el límite de 15MB para envío multimodal directo."
                )
        except Exception as e:
            logger.error(f"Error codificando PDF {filename} a Base64: {e}")

    elif ext == ".docx":
        text = extract_text_from_docx(filepath)
    elif ext == ".pptx":
        text = extract_text_from_pptx(filepath)
    else:
        raise ValueError(f"Extensión de archivo no soportada: {ext}")

    abstract = await generate_document_abstract(text, filename, pdf_base64)

    # Simular la estructura de paper científico
    return {
        "title": filename,
        "authors": "Usuario (Documento Subido)",
        "journal": "Documento Adjunto Local",
        "year": 2026,
        "doi": f"user_upload_{uuid_hash(filename)}",
        "abstract": abstract,
        "extracted_images": extracted_imgs,
    }


async def extract_and_analyze_images(pdf_path: str, run_id: str) -> list:
    """
    Extrae imágenes del PDF usando PyMuPDF (fitz), las guarda en el directorio
    de descargas y las analiza usando Gemini Vision para catalogarlas y captionarlas.
    """
    images_metadata = []
    if not run_id:
        return images_metadata

    out_dir = f"static/downloads/{run_id}/extracted_images"
    os.makedirs(out_dir, exist_ok=True)

    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        logger.error(f"Error abriendo PDF con PyMuPDF: {e}")
        return images_metadata

    image_count = 0
    max_images = 10  # Para no agotar recursos

    for page_num in range(len(doc)):
        if image_count >= max_images:
            break

        page = doc[page_num]
        image_list = page.get_images(full=True)

        for img_idx, img in enumerate(image_list):
            if image_count >= max_images:
                break

            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]

            # Evitar decoraciones pequeñas
            if len(image_bytes) < 10240:  # 10 KB
                continue

            image_filename = f"fig_page_{page_num + 1}_{img_idx + 1}.{image_ext}"
            image_path = os.path.join(out_dir, image_filename)

            with open(image_path, "wb") as f:
                f.write(image_bytes)

            # Codificar imagen a base64 para Gemini Vision
            img_b64 = base64.b64encode(image_bytes).decode("utf-8")
            mime_type = f"image/{image_ext}"
            if image_ext == "jpg":
                mime_type = "image/jpeg"

            # Llamar a Gemini Vision
            prompt = f"""
            Actúa como un experto en imageneología y estadística clínica en cirugía pediátrica.
            Analiza esta imagen extraída de un artículo científico (página {page_num + 1}).

            Tu tarea:
            1. Determina si la imagen es clínicamente relevante: ¿es un gráfico estadístico (barras, líneas, curvas de supervivencia, Forest Plot, etc.), un diagrama de flujo de pacientes (PRISMA, consorte), o una imagen diagnóstica/quirúrgica?
            2. Si NO es relevante (por ejemplo, es un logo, un banner publicitario, una foto decorativa, o un icono simple), responde únicamente con "INVALIDO".
            3. Si SÍ es relevante, genera:
               - Un título científico conciso para la figura (ej: "Figura 1: Curva de supervivencia de Kaplan-Meier...").
               - Una descripción detallada en español (3-4 líneas) interpretando los datos o hallazgos clínicos clave representados en la imagen.

            Tu respuesta debe estar estrictamente en formato JSON con la siguiente estructura:
            {{
              "status": "VALIDO" o "INVALIDO",
              "title": "Título de la figura",
              "caption": "Interpretación y descripción en español clínico"
            }}
            """
            try:
                vision_res = await call_gemini(
                    prompt,
                    json_mode=True,
                    temperature=0.1,
                    inline_data={"mimeType": mime_type, "data": img_b64},
                )
                cleaned_res = vision_res.strip()
                if cleaned_res.startswith("```"):
                    cleaned_res = (
                        cleaned_res.replace("```json", "").replace("```", "").strip()
                    )

                analysis = json.loads(cleaned_res)
                if analysis.get("status") == "VALIDO":
                    relative_path = (
                        f"/static/downloads/{run_id}/extracted_images/{image_filename}"
                    )
                    images_metadata.append(
                        {
                            "file_path": image_path,
                            "url": relative_path,
                            "title": analysis.get("title", f"Figura {image_count + 1}"),
                            "caption": analysis.get("caption", ""),
                            "page": page_num + 1,
                        }
                    )
                    image_count += 1
                    logger.info(f"Imagen válida extraída: {image_filename}")
                else:
                    os.remove(image_path)
            except Exception as e:
                logger.error(f"Error analizando imagen {image_filename}: {e}")
                # Mantenerla como fallback simple
                relative_path = (
                    f"/static/downloads/{run_id}/extracted_images/{image_filename}"
                )
                images_metadata.append(
                    {
                        "file_path": image_path,
                        "url": relative_path,
                        "title": f"Figura {image_count + 1} (Extraída)",
                        "caption": "Figura clínica extraída del documento original.",
                        "page": page_num + 1,
                    }
                )
                image_count += 1

    return images_metadata


async def download_pmc_figures(doi: str, run_id: str) -> list:
    """
    Consulta la API de Europe PMC para obtener y descargar figuras asociadas a un paper por su DOI.
    """
    figures = []
    if not run_id or not doi:
        return figures

    out_dir = f"static/downloads/{run_id}/extracted_images"
    os.makedirs(out_dir, exist_ok=True)

    try:
        # 1. Buscar PMCID por DOI
        search_url = "https://www.ebi.ac.uk/europepmc/webservices/rest/search"
        params = {"query": f"doi:{doi}", "format": "json", "resultType": "lite"}
        async with httpx.AsyncClient(timeout=10.0) as client:
            res = await client.get(search_url, params=params)
            res.raise_for_status()
            results = res.json().get("resultList", {}).get("result", [])

        if not results:
            return figures

        pmcid = results[0].get("pmcid")
        if not pmcid:
            return figures

        # 2. Buscar imágenes asociadas a PMCID
        images_url = f"https://www.ebi.ac.uk/europepmc/webservices/rest/{pmcid}/images"
        async with httpx.AsyncClient(timeout=10.0) as client:
            res_img = await client.get(images_url, params={"format": "json"})
            if res_img.status_code != 200:
                return figures
            images_data = res_img.json().get("images", [])

        # 3. Descargar y registrar las imágenes (límite 3 por paper)
        count = 0
        for img in images_data[:3]:
            img_url = img.get("urls", {}).get("large") or img.get("urls", {}).get(
                "medium"
            )
            if not img_url:
                continue

            img_caption = img.get("caption", "Figura de PubMed Central.")
            img_title = img.get("title", f"Figura PMC ({pmcid})")

            async with httpx.AsyncClient(timeout=15.0) as client:
                res_content = await client.get(img_url)
                if res_content.status_code == 200:
                    ext = "png"
                    if "jpeg" in res_content.headers.get("content-type", ""):
                        ext = "jpg"
                    elif "gif" in res_content.headers.get("content-type", ""):
                        ext = "gif"

                    filename = f"pmc_{pmcid}_{count + 1}.{ext}"
                    filepath = os.path.join(out_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(res_content.content)

                    relative_url = (
                        f"/static/downloads/{run_id}/extracted_images/{filename}"
                    )
                    figures.append(
                        {
                            "file_path": filepath,
                            "url": relative_url,
                            "title": img_title,
                            "caption": img_caption,
                            "doi": doi,
                        }
                    )
                    count += 1
                    logger.info(f"Descargada figura PMC {filename} de DOI: {doi}")
    except Exception as e:
        logger.error(f"Error descargando figuras de PMC para DOI {doi}: {e}")

    return figures


def uuid_hash(name: str) -> str:
    """Genera un hash simple de 8 caracteres a partir del nombre del archivo."""
    return hashlib.md5(name.encode("utf-8")).hexdigest()[:8]
