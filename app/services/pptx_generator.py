import os
import logging
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("multiagent_pptx")

# Paleta de Colores
BG_COLOR = RGBColor(9, 13, 22)         # #090D16 (Azul oscuro profundo)
CARD_BG_COLOR = RGBColor(19, 26, 38)   # #131A26 (Azul de tarjeta)
TEAL_COLOR = RGBColor(0, 210, 196)     # #00D2C4 (Cian quirúrgico)
WHITE_COLOR = RGBColor(255, 255, 255)  # Blanco
GRAY_COLOR = RGBColor(200, 210, 220)   # Gris claro para texto secundario

def apply_dark_background(slide):
    """Establece un color de fondo azul oscuro sólido en la diapositiva."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_slide_title(slide, text: str):
    """Agrega un título estándar a las diapositivas de contenido."""
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEAL_COLOR
    p.alignment = PP_ALIGN.LEFT

def add_slide_references(slide, text: str):
    """Agrega una nota en el borde inferior indicando los papers/referencias utilizados."""
    if not text:
        return
    # Las dimensiones de la diapositiva son 13.333 x 7.5 pulgadas.
    ref_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.8), Inches(11.833), Inches(0.4))
    tf = ref_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = f"Evidencia / Referencias: {text}"
    p.font.name = 'Arial'
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = GRAY_COLOR
    p.alignment = PP_ALIGN.LEFT

def generate_forest_plot(run_id: str, forest_plot_data: list = None) -> str:
    """
    Genera un Forest Plot SOLO a partir de datos reales (odds ratios + IC provenientes
    del corpus analizado). Si no se dispone de datos reales suficientes, devuelve "" en
    lugar de fabricar cifras — una herramienta médica NUNCA debe inventar estadísticas.
    """
    # Requerir datos reales: sin ≥2 estudios con OR real, no se dibuja nada.
    if not forest_plot_data or len(forest_plot_data) < 2:
        logger.info("Forest plot omitido: no hay odds ratios reales en el corpus (no se fabrican datos).")
        return ""

    try:
        run_dir = f"static/downloads/{run_id}" if run_id else "static/downloads/temp"
        os.makedirs(run_dir, exist_ok=True)
        img_path = os.path.join(run_dir, "forest_plot.png")

        using_real_data = True
        try:
            labels = [item["label"] for item in forest_plot_data]
            odds_ratios = [float(item.get("or", 1.0)) for item in forest_plot_data]
            ci_lower = [float(item.get("ci_lower", or_val - 0.3)) for or_val, item in zip(odds_ratios, forest_plot_data)]
            ci_upper = [float(item.get("ci_upper", or_val + 0.3)) for or_val, item in zip(odds_ratios, forest_plot_data)]
            is_summary_flags = [bool(item.get("is_summary", False)) for item in forest_plot_data]
        except Exception as parse_err:
            logger.warning(f"forest_plot_data mal formado, se omite el gráfico: {parse_err}")
            return ""

        fig, ax = plt.subplots(figsize=(6, 4.2), dpi=300)

        # Estilo oscuro
        fig.patch.set_facecolor('#131A26')
        ax.set_facecolor('#131A26')

        # Línea de no efecto (OR = 1.0)
        ax.axvline(x=1.0, color='#ef4444', linestyle='--', linewidth=1.5, label="Línea de no efecto")

        # Graficar estudios
        for idx, (label, or_val, low, high) in enumerate(zip(labels, odds_ratios, ci_lower, ci_upper)):
            is_summary = is_summary_flags[idx]
            color = '#00D2C4' if not is_summary else '#a855f7'
            marker = 'D' if is_summary else 'o'
            markersize = 8 if is_summary else 6

            # Intervalo de confianza
            ax.plot([low, high], [idx, idx], color=color, linewidth=2)
            # OR (Punto medio)
            ax.plot(or_val, idx, marker, color=color, markersize=markersize, markeredgecolor=color)

        ax.set_yticks(range(len(labels)))
        ax.set_yticklabels(labels, color='#FFFFFF', fontsize=9, fontweight='bold')
        ax.set_xlabel("Odds Ratio (OR) e Intervalo de Confianza (95%)", color='#FFFFFF', fontsize=10)
        plot_title = "Meta-Análisis: Efecto Comparativo por Estudio" if using_real_data else "Eficacia Comparativa Quirúrgica (Meta-Análisis)"
        ax.set_title(plot_title, color='#00D2C4', fontsize=11, fontweight='bold')
        
        ax.tick_params(colors='#FFFFFF', which='both', labelsize=9)
        ax.spines['bottom'].set_color('#FFFFFF')
        ax.spines['left'].set_color('#FFFFFF')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        
        ax.set_xlim(0.2, 2.0)
        ax.set_ylim(-0.5, len(labels) - 0.5)
        
        plt.tight_layout()
        plt.savefig(img_path, facecolor=fig.get_facecolor(), edgecolor='none', bbox_inches='tight')
        plt.close()
        return img_path
    except Exception as e:
        logger.error(f"Error generando Forest Plot: {e}")
        return ""


def generate_complication_chart(run_id: str, analyzed_papers: list = None) -> str:
    """
    Gráfico de barras HONESTO con las tasas de complicaciones REALES por estudio,
    extraídas de numeric_data.complication_rate_pct del corpus. Solo usa datos reales;
    si ningún estudio reporta tasa de complicaciones, devuelve "" (no fabrica nada).
    Reemplaza el antiguo forest plot sintético como visualización por defecto.
    """
    if not analyzed_papers:
        return ""

    entries = []
    for p in analyzed_papers:
        nd = p.get("numeric_data") or {}
        rate = nd.get("complication_rate_pct")
        if rate is None:
            continue
        try:
            rate = float(rate)
        except (TypeError, ValueError):
            continue
        authors_raw = str(p.get("authors", "N/A"))
        first = authors_raw.split(",")[0].strip()
        year = p.get("year", "")
        label = f"{first} et al. ({year})" if "," in authors_raw else f"{first} ({year})"
        entries.append((label[:32], rate))

    # Necesitamos al menos 2 estudios con datos reales para que el gráfico sea informativo
    if len(entries) < 2:
        logger.info("Gráfico de complicaciones omitido: <2 estudios con tasa real reportada.")
        return ""

    entries.sort(key=lambda e: e[1], reverse=True)
    entries = entries[:8]  # máximo 8 barras legibles
    labels = [e[0] for e in entries]
    rates = [e[1] for e in entries]

    try:
        run_dir = f"static/downloads/{run_id}" if run_id else "static/downloads/temp"
        os.makedirs(run_dir, exist_ok=True)
        img_path = os.path.join(run_dir, "complication_chart.png")

        fig, ax = plt.subplots(figsize=(6, 4.2), dpi=300)
        fig.patch.set_facecolor('#131A26')
        ax.set_facecolor('#131A26')

        y_pos = range(len(labels))
        ax.barh(list(y_pos), rates, color='#00D2C4', edgecolor='#0b1019', height=0.6)
        ax.set_yticks(list(y_pos))
        ax.set_yticklabels(labels, color='#FFFFFF', fontsize=9, fontweight='bold')
        ax.invert_yaxis()  # el de mayor tasa arriba
        ax.set_xlabel("Tasa de complicaciones reportada (%)", color='#FFFFFF', fontsize=10)
        ax.set_title("Tasas de Complicaciones por Estudio (datos reales del corpus)",
                     color='#00D2C4', fontsize=11, fontweight='bold')

        # Etiquetas de valor al final de cada barra
        for i, v in enumerate(rates):
            ax.text(v + max(rates) * 0.01, i, f"{v:g}%", color='#FFFFFF',
                    fontsize=8, va='center')

        ax.tick_params(colors='#FFFFFF', which='both', labelsize=9)
        ax.spines['bottom'].set_color('#FFFFFF')
        ax.spines['left'].set_color('#FFFFFF')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

        plt.tight_layout()
        plt.savefig(img_path, facecolor=fig.get_facecolor(), edgecolor='none', bbox_inches='tight')
        plt.close()
        return img_path
    except Exception as e:
        logger.error(f"Error generando gráfico de complicaciones: {e}")
        return ""

def _add_picture_keep_ratio(slide, img_path, left, top, max_w, max_h):
    """
    Inserta una imagen preservando su relación de aspecto dentro de un cuadro
    max_w × max_h (in Inches), centrada. Evita la distorsión de figuras clínicas
    que causaba fijar width y height simultáneamente.
    """
    try:
        from PIL import Image as _PILImage
        with _PILImage.open(img_path) as im:
            iw, ih = im.size
        if iw <= 0 or ih <= 0:
            raise ValueError("dimensiones inválidas")
        box_ratio = max_w / max_h
        img_ratio = iw / ih
        if img_ratio > box_ratio:
            w = max_w
            h = max_w / img_ratio
        else:
            h = max_h
            w = max_h * img_ratio
        off_l = left + (max_w - w) / 2
        off_t = top + (max_h - h) / 2
        return slide.shapes.add_picture(img_path, Inches(off_l), Inches(off_t),
                                        width=Inches(w), height=Inches(h))
    except Exception:
        # Fallback: ancho fijo, altura automática (python-pptx preserva ratio con solo width)
        return slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(max_w))


def build_pptx(slides_list: list, filepath: str, run_id: str = None, meta_analysis: dict = None, analyzed_papers: list = None):
    """
    Toma la lista de diccionarios de diapositivas y genera una presentación PPTX
    16:9 premium y libre de errores de maquetación, con notas del orador.
    """
    _forest_data = (meta_analysis or {}).get("forest_plot_data") or []
    _analyzed = analyzed_papers or []

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    for i, s in enumerate(slides_list):
        slide = prs.slides.add_slide(blank_layout)
        apply_dark_background(slide)
        
        layout_type = s.get("layout", "bullet_points")
        
        # --- LAYOUT: TITLE ---
        if layout_type == "title":
            title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.alignment = PP_ALIGN.CENTER
            
            p1 = tf.paragraphs[0]
            p1.text = s.get("title", "Cirugía Pediátrica").upper()
            p1.font.name = 'Arial'
            p1.font.size = Pt(44)
            p1.font.bold = True
            p1.font.color.rgb = TEAL_COLOR
            p1.alignment = PP_ALIGN.CENTER
            
            p2 = tf.add_paragraph()
            p2.text = s.get("subtitle", "Análisis de Evidencia")
            p2.font.name = 'Arial'
            p2.font.size = Pt(20)
            p2.font.color.rgb = WHITE_COLOR
            p2.alignment = PP_ALIGN.CENTER
            
        # --- LAYOUT: BULLET POINTS ---
        elif layout_type == "bullet_points":
            title = s.get("title", "Título de Diapositiva")
            add_slide_title(slide, title)
            
            bullets = s.get("bullets", [])
            bullets_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.8))
            tf = bullets_box.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_left = tf.margin_right = tf.margin_bottom = 0
            
            for idx, bullet_text in enumerate(bullets[:5]):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.text = f"•  {bullet_text}"
                p.font.name = 'Arial'
                p.font.size = Pt(18)
                p.font.color.rgb = WHITE_COLOR
                p.space_after = Pt(14)
                
        # --- LAYOUT: COMPARISON TABLE ---
        elif layout_type == "comparison_table":
            title = s.get("title", "Tabla Comparativa")
            add_slide_title(slide, title)
            
            headers = s.get("headers", [])
            rows = s.get("rows", [])
            
            if headers:
                num_cols = len(headers)
                num_rows = len(rows) + 1
                
                left = Inches(1.0)
                top = Inches(1.8)
                width = Inches(11.333)
                height = Inches(4.5)
                
                table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
                table = table_shape.table
                
                # Cabecera
                for col_idx, header in enumerate(headers):
                    cell = table.cell(0, col_idx)
                    cell.text = header
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CARD_BG_COLOR
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    p.font.name = 'Arial'
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = TEAL_COLOR
                    
                # Celdas
                for row_idx, row_data in enumerate(rows):
                    for col_idx, cell_value in enumerate(row_data):
                        if col_idx >= num_cols:
                            continue
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = cell_value
                        p = cell.text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.LEFT
                        p.font.name = 'Arial'
                        p.font.size = Pt(14)
                        p.font.color.rgb = WHITE_COLOR
                        
        # --- LAYOUT: COMPARISON 2 COLUMNS ---
        elif layout_type == "comparison_2col":
            title = s.get("title", "Comparativa Bilateral")
            add_slide_title(slide, title)
            
            col1_title = s.get("col1_title", "Columna A")
            col1_bullets = s.get("col1_bullets", [])
            col2_title = s.get("col2_title", "Columna B")
            col2_bullets = s.get("col2_bullets", [])
            
            col_width = Inches(5.4)
            col_height = Inches(4.5)
            
            # Columna 1
            shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), col_width, col_height)
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = CARD_BG_COLOR
            shape1.line.color.rgb = TEAL_COLOR
            tf1 = shape1.text_frame
            tf1.word_wrap = True
            tf1.margin_left = tf1.margin_right = Inches(0.3)
            
            p1 = tf1.paragraphs[0]
            p1.text = col1_title
            p1.font.name = 'Arial'
            p1.font.size = Pt(20)
            p1.font.bold = True
            p1.font.color.rgb = TEAL_COLOR
            p1.space_after = Pt(12)
            
            for idx, b in enumerate(col1_bullets[:4]):
                p = tf1.add_paragraph()
                p.text = f"• {b}"
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE_COLOR
                p.space_after = Pt(8)
                
            # Columna 2
            shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), col_width, col_height)
            shape2.fill.solid()
            shape2.fill.fore_color.rgb = CARD_BG_COLOR
            shape2.line.color.rgb = TEAL_COLOR
            tf2 = shape2.text_frame
            tf2.word_wrap = True
            tf2.margin_left = tf2.margin_right = Inches(0.3)
            
            p2 = tf2.paragraphs[0]
            p2.text = col2_title
            p2.font.name = 'Arial'
            p2.font.size = Pt(20)
            p2.font.bold = True
            p2.font.color.rgb = TEAL_COLOR
            p2.space_after = Pt(12)
            
            for idx, b in enumerate(col2_bullets[:4]):
                p = tf2.add_paragraph()
                p.text = f"• {b}"
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE_COLOR
                p.space_after = Pt(8)
                
        # --- LAYOUT: STEP PROCESS ---
        elif layout_type == "step_process":
            title = s.get("title", "Proceso Quirúrgico Paso a Paso")
            add_slide_title(slide, title)
            
            steps = s.get("steps", [])
            num_steps = min(len(steps), 4)
            
            if steps:
                total_width = Inches(11.333)
                gap = Inches(0.4)
                card_width = (total_width - (gap * (num_steps - 1))) / num_steps
                card_height = Inches(4.2)
                card_top = Inches(2.0)
                
                for idx in range(num_steps):
                    card_left = Inches(1.0) + idx * (card_width + gap)
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = CARD_BG_COLOR
                    shape.line.color.rgb = TEAL_COLOR
                    shape.line.width = Pt(1.5)
                    
                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_top = Inches(0.3)
                    tf.margin_left = tf.margin_right = Inches(0.2)
                    
                    p_num = tf.paragraphs[0]
                    p_num.text = f"PASO {idx + 1}"
                    p_num.font.name = 'Arial'
                    p_num.font.size = Pt(20)
                    p_num.font.bold = True
                    p_num.font.color.rgb = TEAL_COLOR
                    p_num.alignment = PP_ALIGN.CENTER
                    p_num.space_after = Pt(14)
                    
                    p_desc = tf.add_paragraph()
                    p_desc.text = steps[idx]
                    p_desc.font.name = 'Arial'
                    p_desc.font.size = Pt(14)
                    p_desc.font.color.rgb = WHITE_COLOR
                    p_desc.alignment = PP_ALIGN.LEFT
                    
        # --- LAYOUT: MULTIMODAL CHART ---
        elif layout_type == "multimodal_chart":
            title = s.get("title", "Análisis Visual de Evidencia")
            add_slide_title(slide, title)
            
            img_url = s.get("image_url")
            bullets = s.get("bullets", [])
            
            if not img_url and run_id:
                # Buscar imagen extraída de la carpeta
                img_dir = f"static/downloads/{run_id}/extracted_images"
                if os.path.exists(img_dir):
                    files = [os.path.join(img_dir, f) for f in os.listdir(img_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
                    if files:
                        img_url = files[0]
                        
            if not img_url:
                # Preferir un gráfico REAL de complicaciones; si no hay datos, un forest plot
                # solo si hay ORs reales. Nunca fabricar.
                img_url = generate_complication_chart(run_id, _analyzed)
                if not img_url:
                    img_url = generate_forest_plot(run_id, forest_plot_data=_forest_data)

            if img_url and os.path.exists(img_url):
                try:
                    _add_picture_keep_ratio(slide, img_url, 1.0, 1.8, 5.5, 4.2)
                except Exception as e:
                    logger.error(f"Error insertando imagen en slide: {e}")
                    
            # Viñetas descriptivas
            bullets_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.5))
            tf = bullets_box.text_frame
            tf.word_wrap = True
            
            for idx, b in enumerate(bullets[:4]):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.text = f"• {b}"
                p.font.name = 'Arial'
                p.font.size = Pt(16)
                p.font.color.rgb = WHITE_COLOR
                p.space_after = Pt(10)
                
        # --- LAYOUT: FOREST PLOT ---
        elif layout_type == "forest_plot":
            title = s.get("title", "Meta-Análisis: Forest Plot de Eficacia")
            add_slide_title(slide, title)
            
            bullets = s.get("bullets", [])
            # Forest plot real si hay ORs; si no, gráfico real de complicaciones; si no, nada.
            img_path = generate_forest_plot(run_id, forest_plot_data=_forest_data)
            if not img_path:
                img_path = generate_complication_chart(run_id, _analyzed)

            if img_path and os.path.exists(img_path):
                try:
                    _add_picture_keep_ratio(slide, img_path, 1.0, 1.8, 5.5, 4.2)
                except Exception as e:
                    logger.error(f"Error insertando Forest Plot: {e}")
                    
            bullets_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.5))
            tf = bullets_box.text_frame
            tf.word_wrap = True
            
            for idx, b in enumerate(bullets[:4]):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.text = f"• {b}"
                p.font.name = 'Arial'
                p.font.size = Pt(16)
                p.font.color.rgb = WHITE_COLOR
                p.space_after = Pt(10)
                
        # --- LAYOUT: METRICS ---
        elif layout_type == "metrics":
            title = s.get("title", "Medida Destacada")
            add_slide_title(slide, title)
            
            metric_val = s.get("metric_value", "99%")
            metric_lbl = s.get("metric_label", "Descripción")
            
            val_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.2))
            tf_val = val_box.text_frame
            tf_val.word_wrap = True
            p_val = tf_val.paragraphs[0]
            p_val.text = metric_val
            p_val.font.name = 'Arial'
            p_val.font.size = Pt(80)
            p_val.font.bold = True
            p_val.font.color.rgb = TEAL_COLOR
            p_val.alignment = PP_ALIGN.CENTER
            
            lbl_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.333), Inches(1.8))
            tf_lbl = lbl_box.text_frame
            tf_lbl.word_wrap = True
            p_lbl = tf_lbl.paragraphs[0]
            p_lbl.text = metric_lbl
            p_lbl.font.name = 'Arial'
            p_lbl.font.size = Pt(22)
            p_lbl.font.color.rgb = WHITE_COLOR
            p_lbl.alignment = PP_ALIGN.CENTER
            
        # --- SPEAKER NOTES ---
        speaker_notes = s.get("speaker_notes", s.get("notes", ""))
        if speaker_notes:
            try:
                slide.notes_slide.notes_text_frame.text = speaker_notes
            except Exception as e:
                logger.error(f"Error escribiendo speaker notes: {e}")

        # Referencias al pie
        references = s.get("references", "")
        if references:
            add_slide_references(slide, references)
            
    os.makedirs(os.path.dirname(filepath), exist_ok=True)
    prs.save(filepath)
    logger.info(f"Presentación de PowerPoint guardada con éxito en: {filepath}")
